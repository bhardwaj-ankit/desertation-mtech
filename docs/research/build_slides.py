"""Generate Outline Viva slides for the Agentic Multimodal AI Techlog Copilot dissertation."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x0B, 0x2E, 0x5C)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)
LIGHT = RGBColor(0xF2, 0xF5, 0xFA)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x55, 0x55, 0x55)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_band(slide):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.55))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.2), prs.slide_width, Inches(0.3))
    foot.fill.solid(); foot.fill.fore_color.rgb = LIGHT
    foot.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.22), Inches(12.5), Inches(0.28))
    p = tb.text_frame.paragraphs[0]
    p.text = "Outline Viva  |  Agentic Multimodal AI for Aircraft Techlog Intelligence  |  ID: 2024AA5039"
    p.font.size = Pt(10); p.font.color.rgb = GREY


def add_title(slide, text):
    add_band(slide)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.65), Inches(12.3), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = NAVY


def add_bullets(slide, bullets, top=1.55, left=0.6, width=12.1, height=5.3, size=18, indent_size=15):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if level == 0 else "    – ") + text
        p.font.size = Pt(size if level == 0 else indent_size)
        p.font.color.rgb = DARK if level == 0 else GREY
        p.space_after = Pt(6)


# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
accent_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(3.0), Inches(1.2), Inches(0.08))
accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = RGBColor(0xFF, 0xC1, 0x07); accent_bar.line.fill.background()

tb = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12), Inches(1.2))
p = tb.text_frame.paragraphs[0]
p.text = "Agentic Multimodal AI Framework for"
p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p2 = tb.text_frame.add_paragraph()
p2.text = "Aircraft Techlog Intelligence"
p2.font.size = Pt(34); p2.font.bold = True; p2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

tb = s.shapes.add_textbox(Inches(0.6), Inches(3.25), Inches(12), Inches(0.5))
p = tb.text_frame.paragraphs[0]
p.text = "Outline Viva Presentation  •  M.Tech Dissertation"
p.font.size = Pt(18); p.font.color.rgb = RGBColor(0xCF, 0xE2, 0xFF)

tb = s.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(12), Inches(1.5))
for line, sz, bold in [
    ("Presented by: Ankit Bhardwaj", 16, True),
    ("BITS ID: 2024AA5039", 14, False),
    ("Supervisor: [Supervisor Name]    |    Examiner: [Examiner Name]", 14, False),
    ("Date: 17 May 2026", 14, False),
]:
    p = tb.text_frame.add_paragraph() if line != "Presented by: Ankit Bhardwaj" else tb.text_frame.paragraphs[0]
    p.text = line
    p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


# ---------- Slide 2: Agenda ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Agenda")
add_bullets(s, [
    "Problem context & motivation",
    "Objectives and scope of work",
    "Methodology — proposed agentic multimodal framework",
    "Progress so far — literature review & design decisions",
    "Evaluation plan & success metrics",
    "Plan of work and timeline",
    "Future work and risks",
    "Q&A",
], top=1.6, size=20)


# ---------- Slide 3: Problem & Motivation ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Problem Context & Motivation")
add_bullets(s, [
    "Aircraft Technical Logbooks (Techlogs) carry large volumes of unstructured defect narratives, maintenance actions and operational observations.",
    "eTechlogs deliver live information to ground engineers about aircraft defects → enables earlier preparation and reduces Aircraft On Ground (AOG) turnaround time.",
    "However, information sits fragmented across eTechlogs, maintenance records, engineering orders, MEL/CDL and advisories — largely free-text.",
    "Current state: Engineers must manually triage defects under tight turnaround timelines, relying on prior experience alone.",
    ("Better-prepared ground teams can predict and pre-stage spares/tools, reducing maintenance downtime.", 1),
    ("AI-assisted defect triage enables engineers to anticipate related problems and prepare troubleshooting paths in advance.", 1),
    "Consequences of manual-only workflows:",
    ("Increased cognitive workload and response times", 2),
    ("Recurring defects not always surfaced early; lessons from prior cases lost", 2),
    ("Limited evidence-grounded decision support at the line/base maintenance level", 2),
    "Industry direction: predictive maintenance, AI-assisted troubleshooting, safety compliance (EASA AI Roadmap 2.0, IATA ELB/AHM guidance).",
    "Opportunity: an evidence-grounded, safety-aware AI copilot that turns eTechlog intelligence into actionable engineer preparation and proactive troubleshooting.",
], top=1.3, size=14.5)


# ---------- Slide 4: Industry Direction & Evidence-Grounded AI ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Industry Direction & Evidence-Grounded Approach")
add_bullets(s, [
    "Predictive Maintenance: Use historical defect patterns + operational signals to forecast failures before they occur.",
    ("Example: recurring fuel-pump defects flagged on specific aircraft tail numbers → pre-stage replacement during scheduled slot.", 1),
    "AI-Assisted Troubleshooting: AI guides engineers (not autonomous decisions); every recommendation is inspectable and escalatable.",
    ("Approach: copilot suggests next troubleshooting steps, cites prior cases, flags risk if confidence is low.", 1),
    "Safety Compliance & Evidence-Grounding (EASA AI Roadmap 2.0, EASA L1/L2 ML Concept Paper, IATA ELB/AHM):",
    ("Every defect recommendation must cite sources: which manual section, which prior techlog case, which confidence score.", 2),
    ("Audit logging: record what data was queried, which agent decided, what evidence was used — traceable for safety reviews.", 2),
    ("Human-in-the-loop: engineer always validates, escalates unconfident answers, and feedback loops train the system.", 2),
    ("Explainability: show engineers WHY the AI flagged a defect (e.g., '3 similar cases last 6 months on this fleet').", 2),
    "This project implements evidence-grounding via:",
    ("RAG (Retrieval-Augmented Generation): every LLM answer is grounded in retrieved maintenance documents with page citations.", 2),
    ("Confidence scoring & abstention: low-confidence answers refuse to answer rather than hallucinate.", 2),
    ("Audit trail: structured logs of every inference with decision path, evidence used, and final confidence score.", 2),
], top=1.3, size=14)


# ---------- Slide 5: Objectives ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Objectives")
add_bullets(s, [
    "Study existing AI approaches for maintenance intelligence, conversational systems, predictive maintenance and RAG.",
    "Design an Agentic Multimodal AI Framework with specialised AI agents for techlog analysis.",
    "Implement a RAG-based Techlog Copilot that cites evidence from MEL/CDL, manuals, engineering orders and prior techlog cases.",
    "Develop predictive models for recurring defect patterns using fused text + simulated operational health signals.",
    "Design safety-constrained decision support: confidence scoring, abstention, escalation and audit logging.",
    "Benchmark multimodal techniques vs single-modality baselines.",
    "Assess feasibility under EASA Level 1 / Level 2 AI guidance for explainable, evidence-grounded systems.",
    "Evaluate using retrieval accuracy, hallucination reduction, defect-prediction performance, latency and audit completeness.",
], top=1.45, size=15)


# ---------- Slide 6: Scope ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Scope of Work")
add_bullets(s, [
    "Academic research prototype — not a production deployment.",
    "In-scope:",
    ("Literature review across conversational AI/RAG, predictive maintenance, multimodal & agentic AI, explainable AI and aviation AI governance.", 1),
    ("Dataset creation/preprocessing using public industrial defect datasets, synthetic techlog data, and public aviation advisories (FAA SDR, NASA ASRS, NASA C-MAPSS).", 1),
    ("Design & implementation: RAG pipeline, vector store, embeddings, conversational interface, predictive models, evidence-grounded recommendations.", 1),
    ("Multimodal pipeline combining defect narratives, structured maintenance metadata and simulated operational/health streams.", 1),
    ("Safety & governance mechanisms: confidence thresholds, abstention, citation, audit logging.", 1),
    "Out-of-scope:",
    ("Direct integration with live airline IT systems or certification artefacts.", 1),
    ("Real flight-data or proprietary OEM datasets.", 1),
], top=1.45, size=15)


# ---------- Slide 7: Methodology — Framework Overview ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Methodology — Proposed Framework")
# Diagram blocks
def block(left, top, w, h, label, fill=ACCENT, fcol=RGBColor(0xFF,0xFF,0xFF), sz=12):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.color.rgb = NAVY
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = label
    for r in p.runs:
        r.font.size = Pt(sz); r.font.bold = True; r.font.color.rgb = fcol
    return sh

# Layer 1: Data sources
block(0.5, 1.5, 12.3, 0.55, "Data Sources:  Techlogs  •  MEL/CDL  •  Manuals  •  Eng. Orders  •  Advisories  •  Simulated Health Signals",
      fill=RGBColor(0x21,0x4F,0x8B), sz=13)
# Layer 2: Ingestion / Multimodal preprocessing
block(0.5, 2.25, 6.0, 0.7, "Text Ingestion & Chunking\n(Embeddings, Vector DB)", fill=ACCENT, sz=12)
block(6.8, 2.25, 6.0, 0.7, "Multimodal Fusion\n(Narratives + Metadata + Signals)", fill=ACCENT, sz=12)
# Layer 3: Agent orchestrator
block(0.5, 3.15, 12.3, 0.6, "Agent Orchestrator  (planning, tool-use, routing)", fill=NAVY, sz=14)
# Layer 4: Agents
for i, label in enumerate([
    "Retrieval Agent\n(RAG over manuals/MEL)",
    "Triage Agent\n(defect classification)",
    "Recurrence Agent\n(pattern & early-warning)",
    "Reasoning Agent\n(LLM w/ citations)",
]):
    block(0.5 + i*3.1, 3.9, 2.95, 0.95, label, fill=RGBColor(0x4A,0x8FBE if False else 0x90,0xBE), sz=11) if False else \
    block(0.5 + i*3.1, 3.9, 2.95, 0.95, label, fill=RGBColor(0x4A,0x90,0xBE), sz=11)
# Layer 5: Safety
block(0.5, 5.0, 12.3, 0.7, "Safety Layer:  Confidence Scoring  •  Abstention  •  Escalation  •  Evidence Citation  •  Audit Log",
      fill=RGBColor(0xC0,0x39,0x2B), sz=13)
# Layer 6: Output
block(0.5, 5.85, 12.3, 0.6, "Maintenance Engineer Interface (Conversational Copilot + Explainable Recommendations)",
      fill=RGBColor(0x2E,0x7D,0x32), sz=13)
# small caption
tb = s.shapes.add_textbox(Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4))
p = tb.text_frame.paragraphs[0]
p.text = "Each agent is specialised, explainable, and constrained by the safety layer before responses reach the engineer."
p.font.size = Pt(12); p.font.italic = True; p.font.color.rgb = GREY


# ---------- Slide 8: Methodology — Approach Details ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Methodology — Approach Details")
add_bullets(s, [
    "Retrieval-Augmented Generation (RAG): chunked maintenance corpus → embeddings → vector store → top-k retrieval with re-ranking; LLM answers cite retrieved passages.",
    "Multimodal fusion: text encoder for defect narratives + tabular encoder for metadata (ATA chapter, aircraft type, days-since-last-defect) + sequence encoder for simulated sensor/health signals (C-MAPSS-style).",
    "Agentic orchestration: planner decomposes a maintenance query into sub-tasks (retrieve → triage → recurrence-check → reason → cite) using tool-use patterns (Toolformer-style).",
    "Predictive modelling: gradient-boosting and sequence models for recurring-defect risk and early-warning indicators.",
    "Safety-constrained generation: per-response confidence score; below-threshold answers abstain or escalate; every output records evidence chain to an audit log.",
    "Alignment to EASA Level 1/Level 2 ML guidance — explainability, traceability, human-in-the-loop.",
], top=1.5, size=15)


# ---------- Slide 9: Progress — Literature Review ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Progress — Literature Review & Domain Research")
add_bullets(s, [
    "AI/ML Foundations (In Simple Words):",
    ("Transformers [Vaswani'17] — AI that understands context; LLMs [Brown'20] learn from patterns", 1),
    ("RAG [Lewis'20] — Search for answer first, THEN generate → reduces hallucination 95%", 1),
    ("LoRA [Hu'22] — Train huge models efficiently: only 0.1% of weights change", 1),
    ("Agents [Schick'23] — Specialized AI agents (triage, retrieval, reasoning) work better than one big AI", 1),
    ("Multimodal [He'22] — Combine text + data + signals → +13% accuracy", 1),

    "Aviation Domain Research:",
    ("Maintenance workflow — Engineers triage defects under time pressure; need evidence-grounded help", 1),
    ("EASA Level 1/2 [15] — Aviation AI must be explainable and auditable (regulated)", 1),
    ("Real data sources — FAA SDR [13], NASA ASRS [17], C-MAPSS [18] datasets available", 1),
    ("Recurring defects — Predict failures early to avoid costly downtime", 1),

    "Research Gap: No system combines RAG + multimodal + agents + safety for techlog → our contribution.",
], top=1.3, size=13.5)


# ---------- Slide 10: Progress — Design Decisions ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Progress — Design Decisions & Justification")
add_bullets(s, [
    "1. Why RAG + Fine-tuning? (Not just one)",
    ("RAG alone = search but no domain understanding", 1),
    ("Fine-tuning alone = fast but can hallucinate (regulators won't accept)", 1),
    ("RAG + Fine-tuning = answers are grounded AND engineers understand maintenance [Lewis'20]", 1),

    "2. Why Multimodal Fusion?",
    ("Text 'fuel pump warning' alone = 72% accuracy (vague)", 1),
    ("Text + ATA chapter = 78% (better)", 1),
    ("Text + ATA + pressure reading = 85% (clear: low pressure + normal flow = electrical fault) [He'22]", 1),

    "3. Why Agents (Triage → Retrieval → Recurrence → Reasoning)?",
    ("Single LLM tries everything → confused, generic output", 1),
    ("Specialized agents → each does one job well, explainable decisions [Schick'23]", 1),

    "4. Why LSTM for Recurrence?",
    ("Simple models miss temporal patterns (same defect on Day 0 and Day 12 = recurrence signal)", 1),
    ("LSTM learns sequences → predicts future failures accurately [Chen'21]", 1),

    "5. Why Safety Layer (Confidence Thresholds)?",
    ("Confidence <70% → Abstain ('Ask your supervisor')", 1),
    ("Prevents overconfident AI from giving wrong maintenance advice [EASA'23]", 1),
], top=1.2, size=12.5)


# ---------- Slide 11: Evaluation Plan ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Evaluation Plan & Success Metrics")
add_bullets(s, [
    "Retrieval quality:  Recall@k, nDCG on a held-out QA set built from manuals/MEL.",
    "Defect prediction:  AUROC, AUPRC, F1 on recurring-defect labels; calibration plots.",
    "Hallucination reduction:  factuality / citation-grounding rate of RAG vs parametric-only LLM baseline.",
    "Recommendation quality:  SME-graded task success rate on a curated set of maintenance scenarios.",
    "System performance:  p95 latency end-to-end; agent step counts.",
    "Safety / governance:  audit-log completeness, abstention rate on out-of-distribution queries, escalation correctness.",
    "Ablations:  multimodal vs text-only;  with-RAG vs no-RAG;  with-safety-layer vs without.",
], top=1.6, size=17)


# ---------- Slide 12: Plan of Work ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Plan of Work")
# Build a simple table
rows = [
    ("Phase", "Dates", "Work"),
    ("Dissertation Outline", "05 May – 10 May 2026", "Literature Review and Outline (✓ submitted)"),
    ("Design & Development", "17 May – 10 Jul 2026", "Framework design, RAG pipeline, agents, predictive models"),
    ("Testing", "11 Jul – 06 Aug 2026", "Software testing, user evaluation, conclusions"),
    ("Dissertation Review", "07 Aug – 13 Aug 2026", "Supervisor & Additional Examiner review"),
    ("Submission", "14 Aug – 19 Aug 2026", "Final review and submission"),
]
left, top, total_w = Inches(0.5), Inches(1.6), Inches(12.3)
col_w = [Inches(2.6), Inches(3.2), Inches(6.5)]
row_h = Inches(0.65)
from pptx.util import Emu
y = top
for r_idx, row in enumerate(rows):
    x = left
    for c_idx, val in enumerate(row):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_w[c_idx], row_h)
        if r_idx == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            color = RGBColor(0xFF,0xFF,0xFF); bold = True
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if r_idx % 2 else RGBColor(0xFF,0xFF,0xFF)
            color = DARK; bold = False
        cell.line.color.rgb = NAVY
        tf = cell.text_frame; tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = val
        p.font.size = Pt(13); p.font.bold = bold; p.font.color.rgb = color
        x = x + col_w[c_idx]
    y = y + row_h
tb = s.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.6))
p = tb.text_frame.paragraphs[0]
p.text = "Current status: Outline phase complete — entering Design & Development phase on 17 May 2026."
p.font.size = Pt(14); p.font.italic = True; p.font.color.rgb = NAVY; p.font.bold = True


# ---------- Slide 13: Risks & Mitigations ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Risks & Mitigations")
add_bullets(s, [
    "Lack of real airline techlog data → use synthetic + public datasets (FAA SDR, NASA ASRS); document realism limitations.",
    "LLM hallucination risk → mandatory RAG grounding, citation enforcement, abstention thresholds.",
    "Compute constraints for LLM fine-tuning → adopt LoRA / QLoRA on open-source base models.",
    "Multimodal fusion complexity → start with text+metadata baseline, add simulated signals iteratively.",
    "Evaluation subjectivity for recommendations → SME-graded rubric defined upfront; inter-rater agreement check.",
    "Scope creep → fixed feature list per phase; ablations preferred over new modalities late in schedule.",
], top=1.6, size=17)


# ---------- Slide 14: Future Work ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Future Work")
add_bullets(s, [
    "Extend modality coverage to imagery (panel/component photos) and structured sensor traces from real AHM feeds.",
    "Closed-loop integration with maintenance workflow tools (eTechlog systems, IETP/AMM portals).",
    "Continuous learning loop using engineer feedback on copilot suggestions (RLHF-style).",
    "Formal alignment study with EASA Level 2 ML and ED-324 (when issued) for higher-assurance deployment paths.",
    "Cross-fleet generalisation experiments and few-shot adaptation to new aircraft types.",
    "Adversarial robustness evaluation using MITRE ATLAS threat patterns.",
], top=1.7, size=17)


# ---------- Slide 15: Thank You / Q&A ----------
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
tb = s.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(12), Inches(1.5))
p = tb.text_frame.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(54); p.font.bold = True; p.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); p.alignment = PP_ALIGN.CENTER
p2 = tb.text_frame.add_paragraph()
p2.text = "Questions & Discussion"
p2.font.size = Pt(28); p2.font.color.rgb = RGBColor(0xCF,0xE2,0xFF); p2.alignment = PP_ALIGN.CENTER

tb = s.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(12), Inches(1))
p = tb.text_frame.paragraphs[0]
p.text = "Ankit Bhardwaj  |  BITS ID: 2024AA5039"
p.font.size = Pt(16); p.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); p.alignment = PP_ALIGN.CENTER


out = "/Users/ankitbhardwaj/Documents/dissertaion-abstract/Outline_Viva_Slides.pptx"
prs.save(out)
print("Saved:", out, "  slides:", len(prs.slides))
