"""Generate Mid-Semester Viva slides for the Agentic Multimodal AI Techlog Copilot dissertation.

Content sourced from MID_SEM_REPORT_2024AA05939.md and data/dataset_summary.json.
Covers the BITS-mandated sections: Introduction, Literature Review, Research
Methodology, Data Collection / Progress, and faculty-feedback updates.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x0B, 0x2E, 0x5C)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)
LIGHT = RGBColor(0xF2, 0xF5, 0xFA)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x55, 0x55, 0x55)
GOLD = RGBColor(0xFF, 0xC1, 0x07)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_band(slide):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.55))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.2), prs.slide_width, Inches(0.3))
    foot.fill.solid(); foot.fill.fore_color.rgb = LIGHT
    foot.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.22), Inches(12.5), Inches(0.28))
    p = tb.text_frame.paragraphs[0]
    p.text = "Mid-Semester Viva  |  Agentic Multimodal AI for Aircraft Techlog Intelligence  |  Ankit Bhardwaj  |  2024AA05939"
    p.font.size = Pt(10); p.font.color.rgb = GREY


def add_title(slide, text, sub=None):
    add_band(slide)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(12.3), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = NAVY
    if sub:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.52), Inches(1.28), Inches(2.0), Inches(0.045))
        bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background()
        tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.32), Inches(12.3), Inches(0.4))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = sub
        p2.font.size = Pt(14); p2.font.italic = True; p2.font.color.rgb = ACCENT


def add_bullets(slide, bullets, top=1.55, left=0.6, width=12.1, height=5.3, size=18, indent_size=15):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = {0: "• ", 1: "    – ", 2: "        · "}.get(level, "• ")
        p.text = prefix + text
        p.font.size = Pt(size if level == 0 else indent_size)
        p.font.bold = (level == 0)
        p.font.color.rgb = DARK if level == 0 else GREY
        p.space_after = Pt(5)


def block(slide, left, top, w, h, label, fill=ACCENT, fcol=WHITE, sz=12, line=NAVY):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = label
    for r in p.runs:
        r.font.size = Pt(sz); r.font.bold = True; r.font.color.rgb = fcol
    return sh


def make_table(slide, rows, left, top, col_w, row_h=0.5, header_sz=12, body_sz=11):
    from pptx.util import Emu
    y = Inches(top)
    for r_idx, row in enumerate(rows):
        x = Inches(left)
        for c_idx, val in enumerate(row):
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(col_w[c_idx]), Inches(row_h))
            if r_idx == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
                color = WHITE; bold = True; sz = header_sz
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT if r_idx % 2 else WHITE
                color = DARK; bold = False; sz = body_sz
            cell.line.color.rgb = RGBColor(0xC9, 0xD3, 0xE0)
            tf = cell.text_frame
            tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.06)
            tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
            tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.text = str(val)
            p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = color
            if c_idx > 0 and r_idx > 0:
                p.alignment = PP_ALIGN.LEFT
            x = Emu(x + Inches(col_w[c_idx]))
        y = Emu(y + Inches(row_h))


# ========== Slide 1: Title ==========
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), prs.slide_height)
side.fill.solid(); side.fill.fore_color.rgb = GOLD; side.line.fill.background()

tb = s.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(12), Inches(0.5))
p = tb.text_frame.paragraphs[0]
p.text = "BIRLA INSTITUTE OF TECHNOLOGY & SCIENCE, PILANI"
p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = RGBColor(0xCF, 0xE2, 0xFF)
p2 = tb.text_frame.add_paragraph()
p2.text = "M.Tech (Artificial Intelligence & Machine Learning)  •  AIMLCZG628T Dissertation / Project Work"
p2.font.size = Pt(12); p2.font.color.rgb = RGBColor(0xCF, 0xE2, 0xFF)

accent_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(2.05), Inches(1.4), Inches(0.08))
accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = GOLD; accent_bar.line.fill.background()

tb = s.shapes.add_textbox(Inches(0.7), Inches(2.25), Inches(12), Inches(1.7))
p = tb.text_frame.paragraphs[0]
p.text = "An Agentic Multimodal AI Framework for"
p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHITE
p2 = tb.text_frame.add_paragraph()
p2.text = "Aircraft Techlog Intelligence"
p2.font.size = Pt(36); p2.font.bold = True; p2.font.color.rgb = WHITE

tb = s.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(12), Inches(0.5))
p = tb.text_frame.paragraphs[0]
p.text = "Mid-Semester Progress Viva  •  The “Techlog Copilot”"
p.font.size = Pt(19); p.font.color.rgb = GOLD

tb = s.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(12), Inches(1.8))
for line, sz, bold in [
    ("Presented by:  Ankit Bhardwaj   (ID: 2024AA05939)", 17, True),
    ("Dissertation carried out at:  The Emirates Group, Dubai, UAE", 14, False),
    ("Under the Supervision of:  Alla Venkatesh Rao", 14, False),
    ("Birla Institute of Technology & Science, Pilani (Rajasthan)  •  June 2026", 14, False),
]:
    first = line.startswith("Presented")
    p = tb.text_frame.paragraphs[0] if first else tb.text_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = WHITE
    p.space_after = Pt(4)


# ========== Slide 2: Agenda ==========
s = prs.slides.add_slide(BLANK)
add_title(s, "Agenda", "What this viva covers")
add_bullets(s, [
    "Introduction & problem context",
    "Aim, objectives and scope",
    "Literature review — foundations & research gap",
    "Research methodology — agentic multimodal architecture",
    "Data collection & preparation — work completed this phase",
    "Dataset description & statistics",
    "Evaluation plan & success metrics",
    "Project plan, status and next steps",
    "Q & A",
], top=1.7, size=19)


# ========== Slide 3: Introduction & Problem Context ==========
s = prs.slides.add_slide(BLANK)
add_title(s, "Introduction & Problem Context", "Why aircraft techlog intelligence?")
add_bullets(s, [
    "The aircraft Technical Logbook is one of the most operationally critical documents in commercial aviation — every reported defect, maintenance action and deferred (MEL) item is recorded there.",
    "At a large hub carrier (several hundred aircraft), techlog entries run into the thousands per week.",
    "Engineers must triage defects within a 3–5 minute line-maintenance window by cross-referencing:",
    ("Maintenance manuals (AMM), MEL/CDL provisions, Engineering Orders, and prior resolved cases.", 1),
    ("These sit in systems that do not talk to each other — largely free text.", 1),
    "Today's tools are mostly keyword retrieval: no integration across sources, and no predictive layer for recurring faults.",
    "Cost of a missed or delayed resolution is high — in Aircraft-on-Ground (AOG) time and in safety exposure.",
    "Research opportunity: combine Retrieval-Augmented Generation, multimodal deep learning and agentic AI to fit the operational constraints of an airline MRO environment.",
], top=1.7, size=15.5)


# ========== Slide 4: Aim, Objectives & Scope ==========
s = prs.slides.add_slide(BLANK)
add_title(s, "Aim, Objectives & Scope", "The “Techlog Copilot”")
add_bullets(s, [
    "Aim: an Agentic Multimodal AI framework that assists maintenance engineers with evidence-backed, safety-aware decision support.",
    "Objectives:",
    ("Triage defects across 8 ATA-aligned fault categories with confidence scores.", 1),
    ("Retrieve cited evidence from a curated knowledge corpus (RAG).", 1),
    ("Predict recurring defects on a per-tail (per-aircraft) basis within a 30-day horizon.", 1),
    ("Generate explainable recommendations under safety-governance constraints (EASA L1/L2).", 1),
    "Scope: academic research prototype — not a production deployment.",
    ("No proprietary airline data is used; all real data is from public sources, synthetic data is clearly separated throughout.", 1),
], top=1.7, size=16)


# ========== Slide 5: Literature Review - Foundations ==========
s = prs.slides.add_slide(BLANK)
add_title(s, "Literature Review — Foundations", "Peer-reviewed backing for each design choice")
make_table(s, [
    ("Area", "Key Reference", "Insight informing this work"),
    ("Retrieval-Augmented Generation", "Lewis et al., NeurIPS 2020", "Ground answers in retrieved documents → sharply reduces hallucination; every answer is citable."),
    ("Transformer / LLM backbone", "Vaswani et al., 2017; BERT, 2019", "Contextual encoders for noisy free-text defect narratives."),
    ("Parameter-efficient fine-tuning", "LoRA (Hu'22); QLoRA (Dettmers'23)", "Specialise a 7B model by training ~0.1% of weights — feasible without large GPU clusters."),
    ("Tool-using / agentic AI", "Toolformer (Schick et al., 2023)", "Specialised agents outperform a single monolithic LLM and are more explainable."),
    ("Multimodal learning", "He et al., 2022", "Fusing text + metadata + signals improves classification robustness."),
    ("Sequence modelling", "C-MAPSS / PHM (Saxena'08)", "Run-to-failure sensor sequences support recurrence & degradation modelling."),
    ("Aviation AI governance", "EASA L1/L2 ML Concept Paper '23; NIST AI RMF", "Explainability, traceability, abstention and human-in-the-loop are mandatory."),
], left=0.5, top=1.65, col_w=[3.0, 3.3, 6.5], row_h=0.66, header_sz=12, body_sz=11)


# ========== Slide 6: Literature Review - Research Gap ==========
s = prs.slides.add_slide(BLANK)
add_title(s, "Literature Review — Research Gap", "Where this dissertation contributes")
add_bullets(s, [
    "Existing maintenance-support tools are largely retrieval-only: keyword search over PDF manuals, MEL lookup by ATA chapter, recurring faults left to institutional memory.",
    "Individual building blocks are mature in isolation — but no published system combines them for the techlog setting:",
    ("RAG grounding  +  multimodal fusion (text + metadata + signals)  +  agent orchestration  +  an explicit safety / abstention layer.", 1),
    "Aviation imposes constraints most general RAG/agent papers do not address:",
    ("Evidence citation for every recommendation; auditable decision trails; mandatory abstention on safety-critical questions; sub-5-minute turnaround.", 1),
    "Contribution: an integrated, safety-aware agentic framework that fits airline MRO operating constraints, evaluated on real + synthetic aviation data.",
], top=1.7, size=15.5)


# ========== Slide 7: Methodology - Architecture diagram ==========
s = prs.slides.add_slide(BLANK)
add_title(s, "Research Methodology — System Architecture", "Multi-agent pipeline orchestrated with LangGraph")
block(s, 0.5, 1.55, 12.3, 0.5,
      "Engineer Query (natural language)", fill=RGBColor(0x21, 0x4F, 0x8B), sz=13)
block(s, 0.5, 2.18, 12.3, 0.5,
      "Agent Orchestrator  (LangGraph: planning, routing, tool-use)", fill=NAVY, sz=13)
agents = [
    "Triage Agent\nPyTorch MLP\n8-class label + confidence",
    "Recurrence Agent\nPyTorch LSTM\nrepeat-risk ≤ 30 days",
    "Retrieval Agent\nLangChain + Chroma\ntop-k cited documents",
    "Reasoning Agent\nMistral-7B + LoRA\nranked cited recommendation",
]
for i, label in enumerate(agents):
    block(s, 0.5 + i * 3.1, 2.85, 2.95, 1.15, label, fill=RGBColor(0x4A, 0x90, 0xBE), sz=11)
block(s, 0.5, 4.25, 12.3, 0.62,
      "Safety Layer:  confidence gate  •  abstention (MUST_ABSTAIN)  •  human escalation  •  audit log",
      fill=RED, sz=13)
block(s, 0.5, 5.0, 12.3, 0.55,
      "Engineer Interface  (Streamlit UI + FastAPI backend) — explainable, cited recommendations",
      fill=GREEN, sz=13)
tb = s.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.6))
p = tb.text_frame.paragraphs[0]
p.text = ("Each agent is specialised and explainable; every response passes the safety layer "
          "(confidence gating + abstention + audit) before reaching the engineer.")
p.font.size = Pt(12.5); p.font.italic = True; p.font.color.rgb = GREY


# ========== Slide 8: Methodology - Agent details ==========
s = prs.slides.add_slide(BLANK)
add_title(s, "Research Methodology — Component Detail", "How each agent works")
add_bullets(s, [
    "Triage agent — multimodal MLP fusing three streams: defect narrative (sentence-transformer embedding) + structured metadata (ATA chapter, flight phase, aircraft type, severity) + 42 C-MAPSS-style sensor features (last value & slope). Output: 8 fault classes from ATA iSpec 2200 groupings.",
    "Recurrence agent — LSTM over the chronologically ordered prior-defect sequence per tail number; features include category, ATA chapter, severity and days-since-last-event. Binary output: same category recurs within 30 days.",
    "Retrieval agent — RAG pipeline; Chroma vector store + sentence-transformer embeddings over a 1,213-document corpus (real FAA ADs + synthetic AMM/MEL/EO/advisories + de-identified ASRS cases).",
    "Reasoning agent — Mistral-7B with LoRA adapters, fine-tuned on 26,210 real defect→action pairs from FAA SDRs; produces a ranked, cited recommendation.",
    "Safety layer — confidence gate suppresses low-confidence output; 9 MUST_ABSTAIN question types are declined and escalated; every decision (incl. abstentions) is written to an audit log. Aligned to EASA L1/L2 guidance.",
], top=1.7, size=14)


# ========== Slide 9: Data Collection - Real public data ==========
s = prs.slides.add_slide(BLANK)
add_title(s, "Data Collection — Real Public Data", "Four real aviation sources integrated  (zero proprietary data)")
add_bullets(s, [
    "NASA C-MAPSS Turbofan Degradation (FD001–FD004) — multivariate run-to-failure time series (3 operational settings + 21 sensors); 14 train/test/RUL files. Drives the signal modality of the triage classifier.",
    "FAA Service Difficulty Reports (SDR) 2023–2025 — ~197,000 defect reports, 76-column schema. Normalised to 196,118 records; 166,037 mapped to one of 8 classes via ATA lookup + keyword fallback (30,081 fell into OTHER — a documented limitation).",
    ("Regex splitter extracted 26,210 defect→corrective-action pairs from the free-text Discrepancy field — the primary supervised signal for the reasoning agent.", 1),
    "NASA ASRS Report Sets — 3 curated PDFs (maintenance, fuel, cabin fumes/smoke); parsed with pypdf, segmented by ACN → 150 structured CASE documents added to the corpus.",
    "FAA Airworthiness Directives — 500 AD rules via the Federal Register API (49 with full text from govinfo.gov); all 500 included in the RAG corpus.",
], top=1.7, size=14)


# ========== Slide 10: Data Collection - Synthetic data ==========
s = prs.slides.add_slide(BLANK)
add_title(s, "Data Collection — Synthetic Corpus", "Schema-consistent labelled data the public domain lacks")
add_bullets(s, [
    "Synthetic techlogs — 4,425 records across 70 tail numbers (A6- style fleet: A320/A321/A330/A350/B737/B777/B787). Each carries category, ATA chapter, component, severity, flight phase, station, plus raw (shorthand) and clean narratives. Per-tail histories are internally consistent with 30-day recurrence flags.",
    "Knowledge corpus — 563 synthetic documents (42 AMM snippets, 16 MEL entries, 40 ADs, 35 Engineering Orders, 30 advisories, 400 prior cases). Component→ATA pairing validated programmatically (zero mismatches). With 500 real ADs + 150 ASRS cases → 1,213 retrievable documents.",
    "C-MAPSS-style signal windows — 3,915 windows (3 settings + 21 sensors) with defect-specific degradation signatures injected (e.g. pressure drift for hydraulic, temperature rise for bleed air). Fusion-ready last-value & slope features per sensor.",
    "QA benchmark — 189 question-answer pairs across all 8 categories: 163 MUST_CITE, 17 CAN_ANSWER, 9 MUST_ABSTAIN. Used to evaluate Recall@k, nDCG, citation correctness, hallucination rate, and abstention precision/recall.",
    "Pipeline is fully reproducible: deterministic (seed=42), runs end-to-end via  python3 -m src.data.prepare_all.",
], top=1.7, size=13.5)


# ========== Slide 11: Dataset Summary table ==========
s = prs.slides.add_slide(BLANK)
add_title(s, "Dataset Description & Statistics", "Table 1 — Dataset summary")
make_table(s, [
    ("Dataset", "Type", "Volume"),
    ("NASA C-MAPSS FD001–FD004", "Real", "14 files (train / test / RUL per sub-dataset)"),
    ("FAA SDR 2023–2025", "Real", "~197,000 rows, 76 columns"),
    ("NASA ASRS (3 report sets)", "Real (de-identified)", "3 PDFs → 150 parsed reports"),
    ("FAA Airworthiness Directives", "Real", "500 metadata records (49 full-text)"),
    ("Synthetic Techlogs", "Synthetic", "4,425 records, 70 tail histories"),
    ("Knowledge Corpus", "Mixed", "1,213 documents total"),
    ("C-MAPSS-style Signal Windows", "Synthetic", "3,915 windows / 117,450 cycle-rows"),
    ("SDR Normalised Records", "Processed", "196,118 records"),
    ("Defect→Action Pairs", "Processed", "26,210 pairs"),
    ("QA Benchmark", "Processed", "189 items (163 cite / 17 answer / 9 abstain)"),
], left=0.6, top=1.62, col_w=[4.6, 3.0, 4.5], row_h=0.46, header_sz=13, body_sz=11.5)


# ========== Slide 12: Splits ==========
s = prs.slides.add_slide(BLANK)
add_title(s, "Train / Validation / Test Splits", "Chronological splitting prevents temporal leakage")
add_bullets(s, [
    "Splits are chronological (earliest 80% train / next 10% val / final 10% test) — critical for recurrence, where using future events to predict the past would inflate metrics.",
    "Triage task:  3,540 train  /  442 val  /  443 test.",
    ("Training set roughly uniform across the 8 classes (400–480 per class); 3,915 examples carry the signal modality.", 1),
    "Recurrence task:  3,470 train  /  372 val  /  373 test.",
    ("Positive rate (recurrence within 30 days) ≈ 36% in the training set (1,240 / 3,470).", 1),
    "8 triage classes: Electrical Fault, Hydraulic Leak, Fuel System, Structural Damage, Sensor Indication, Pneumatic/Bleed, Landing Gear/Brake, Powerplant.",
], top=1.7, size=15.5)


# ========== Slide 13: Progress snapshot ==========
s = prs.slides.add_slide(BLANK)
add_title(s, "Progress Snapshot — Work Completed", "Mid-semester deliverable: the full data foundation")
left_items = [
    ("DONE", "Literature review & dissertation outline", GREEN),
    ("DONE", "Reproducible data pipeline — src/data/, 11 modules, seed=42", GREEN),
    ("DONE", "4 real public datasets downloaded & normalised", GREEN),
    ("DONE", "Synthetic techlogs, corpus, signals & QA benchmark", GREEN),
    ("DONE", "Train / val / test splits for triage & recurrence", GREEN),
    ("NEXT", "RAG pipeline & vector store — current focus (Phase 3)", ACCENT),
    ("NEXT", "Triage (MLP) & recurrence (LSTM) model training", GREY),
    ("NEXT", "LoRA fine-tuning + agent integration + safety layer", GREY),
    ("NEXT", "FastAPI + Streamlit UI & full system evaluation", GREY),
]
y = 1.62
for tag, text, col in left_items:
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y), Inches(1.1), Inches(0.42))
    chip.fill.solid(); chip.fill.fore_color.rgb = col; chip.line.fill.background()
    tf = chip.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.text = tag
    p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = WHITE
    tb = s.shapes.add_textbox(Inches(1.85), Inches(y), Inches(10.8), Inches(0.42))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(15); p.font.color.rgb = DARK
    y += 0.50
tb = s.shapes.add_textbox(Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.6))
p = tb.text_frame.paragraphs[0]
p.text = ("Mid-semester deliverable = the complete, reproducible data foundation (Phases 1–2). "
          "Modelling & agent code begins next.")
p.font.size = Pt(13.5); p.font.italic = True; p.font.bold = True; p.font.color.rgb = NAVY


# ========== Slide 14: Evaluation Plan ==========
s = prs.slides.add_slide(BLANK)
add_title(s, "Evaluation Plan & Success Metrics", "How the system will be judged")
add_bullets(s, [
    "Retrieval quality:  Recall@k and nDCG on the 189-item QA benchmark.",
    "Defect prediction:  AUROC, AUPRC, F1 on triage & recurrence; calibration plots.",
    "Hallucination / grounding:  citation-correctness & factuality of RAG vs a parametric-only LLM baseline.",
    "Safety / governance:  abstention precision-recall on MUST_ABSTAIN items; audit-log completeness; escalation correctness.",
    "System performance:  end-to-end p95 latency; agent step counts.",
    "Ablations:  multimodal vs text-only  |  with-RAG vs no-RAG  |  with-safety-layer vs without.",
], top=1.7, size=16.5)


# ========== Slide 15: Project Plan / Timeline ==========
s = prs.slides.add_slide(BLANK)
add_title(s, "Project Plan & Status", "Table 2 — from outline to submission")
make_table(s, [
    ("Phase", "Window", "Status"),
    ("1. Dissertation Outline (lit. review + outline)", "05–10 May 2026", "COMPLETED"),
    ("2. Data Acquisition & Preparation", "10 May – 15 Jun 2026", "COMPLETED"),
    ("3. RAG Pipeline & Vector Store", "16 Jun – 05 Jul 2026", "IN PROGRESS"),
    ("4. Triage & Recurrence Models", "06–20 Jul 2026", "PENDING"),
    ("5. LLM Fine-tuning & Agent Integration", "21 Jul – 06 Aug 2026", "PENDING"),
    ("6. API, UI & Evaluation", "07–13 Aug 2026", "PENDING"),
    ("7. Dissertation Review & Submission", "14–19 Aug 2026", "PENDING"),
], left=0.6, top=1.65, col_w=[6.0, 3.6, 2.5], row_h=0.55, header_sz=13, body_sz=12)
tb = s.shapes.add_textbox(Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.6))
p = tb.text_frame.paragraphs[0]
p.text = ("We are here: data foundation complete (Phases 1–2); RAG pipeline now under way (Phase 3).")
p.font.size = Pt(14); p.font.italic = True; p.font.bold = True; p.font.color.rgb = NAVY


# ========== Slide 16: Risks & Known Limitations ==========
s = prs.slides.add_slide(BLANK)
add_title(s, "Risks & Known Limitations", "Identified early, with mitigations")
add_bullets(s, [
    "No real airline techlog data → synthetic + public datasets used; realism limitations documented; synthetic clearly separated.",
    "30,081 SDR records fell into OTHER (heuristic ATA mapping) → preserved & flagged; future ML classifier can recover them.",
    "LLM hallucination → mandatory RAG grounding, citation enforcement, and abstention thresholds.",
    "Compute limits for 7B fine-tuning → LoRA / QLoRA on open-source base models.",
    "Multimodal fusion complexity → start text+metadata baseline, add signals iteratively.",
    "Recommendation-quality evaluation is subjective → SME-graded rubric defined upfront.",
], top=1.7, size=16)


# ========== Slide 17: Selected References ==========
s = prs.slides.add_slide(BLANK)
add_title(s, "Selected References", "Core papers & standards underpinning the framework")
add_bullets(s, [
    "[1]  Lewis, P. et al. “Retrieval-Augmented Generation for Knowledge-Intensive NLP Tasks.” NeurIPS, 2020.",
    "[2]  Vaswani, A. et al. “Attention Is All You Need.” NeurIPS, 2017.",
    "[3]  Hu, E. J. et al. “LoRA: Low-Rank Adaptation of Large Language Models.” ICLR, 2022.",
    "[4]  Schick, T. et al. “Toolformer: Language Models Can Teach Themselves to Use Tools.” NeurIPS, 2023.",
    "[5]  Saxena, A. et al. “Damage Propagation Modeling for Aircraft Engine Run-to-Failure Simulation (C-MAPSS).” PHM, 2008.",
    "[6]  EASA Concept Paper: First Usable Guidance for Level 1 & 2 Machine Learning Applications, Issue 2, 2023.",
    "[7]  NIST AI Risk Management Framework (AI RMF 1.0), 2023.",
    "Data sources:  FAA SDR  •  NASA ASRS  •  NASA C-MAPSS  •  FAA Airworthiness Directives (all public).",
], top=1.7, size=15)


# ========== Slide 18: Thank You ==========
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), prs.slide_height)
side.fill.solid(); side.fill.fore_color.rgb = GOLD; side.line.fill.background()
tb = s.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(12), Inches(1.5))
p = tb.text_frame.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(54); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
p2 = tb.text_frame.add_paragraph()
p2.text = "Questions & Discussion"
p2.font.size = Pt(26); p2.font.color.rgb = GOLD; p2.alignment = PP_ALIGN.CENTER
tb = s.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(12), Inches(1))
p = tb.text_frame.paragraphs[0]
p.text = "Ankit Bhardwaj   |   2024AA05939   |   M.Tech AI & ML, BITS Pilani"
p.font.size = Pt(16); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER


out = "/Users/ankitbhardwaj/Documents/dissertaion-abstract/MidSem_Viva_Slides_2024AA05939.pptx"
prs.save(out)
print("Saved:", out, "  slides:", len(prs.slides))
